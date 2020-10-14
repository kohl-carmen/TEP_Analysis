

from neuron import h
h.nrn_load_dll('C:\\Users\\ckohl\\Miniconda3\\Lib\\site-packages\\hnn_core\\nrnmech.dll')
import os.path as op
from hnn_core import simulate_dipole, Params, Network, read_params
import json
import numpy as np
import matplotlib.pyplot as plt
import sys        
sys.path.append('C:\\Users\\ckohl\\Desktop\\Current\\TMS\\git-TEP_Analysis\\hnn_core_for_TEP')
import my_hnn_core_functions as m
from pptx import Presentation #https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
from pptx.util import Inches
import pickle
import seaborn

my_param_out_dir='C:\\Users\\ckohl\\hnn_out\\'
dump_dir="C:\\Users\\ckohl\\Desktop\\Current\\TMS\\Core output"     

supra=True
#-------------------------------------
#PPT init
#-------------------------------------
left = 0
top=Inches(0.25)
height = Inches(7)
width=Inches(10)
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "TEP Predictions GUI"
subtitle.text = "plot_tep_predictions_from_gui.py"

#-------------------------------------
# get default
#-------------------------------------
if supra:
    dict_param=m.get_dict_param(my_param_out_dir,'supra_Sarah')
else:
    dict_param=m.get_dict_param(my_param_out_dir,'default_Sarah')
default_params=Params(dict_param)
net_default = Network(default_params)
dpls_default = simulate_dipole(net_default, n_trials=1)
# m.core_output_basic(dpls_default,net,name,False,False,True, False)

all_times=[0, 1, 2, 3, 4, 5, 6, 7, 10, 12, 15, 20, 40, 60, 80, 100, 120, 140, 160, 180, 200, 220, 240, 280, 300, 340, 400]

for time_i in all_times:
    if supra:
        file_dir='TEPpred_supra_Sarah_2_'+str(time_i)+'ms'
        name='suprax2 delay:'+str(time_i)+'ms'
    else:
        file_dir='TEPpred_degfault_Sarah_2_'+str(time_i)+'ms'
        name='defaultx2 delay:'+str(time_i)+'ms'
    dict_param=m.get_dict_param(my_param_out_dir,file_dir)
    these_params=[]
    these_params=Params(dict_param)
    net=[]
    dpls=[]
    net = Network(these_params)
    dpls = simulate_dipole(net, n_trials=1)   
    # m.core_output_basic(dpls,net,name,False,False,True, False)
    #m.plot_hnn_core_output(dpls,net,p,True,False,True,param_oi_t,these_params,base_params)
    # plt.pause(.1)
    # plt.savefig(dump_dir+'\\temp.png')
    # img_path = dump_dir+'\\temp.png'
    # blank_slide_layout = prs.slide_layouts[6]
    # slide = prs.slides.add_slide(blank_slide_layout)
    # pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
    # plt.close()
    #figure:
    fig = plt.figure(figsize=(20,10))
    plt.ion()
    if supra:
        ylms=[-210,200]
    else:
        ylms=[-150,100]
    xlim=net.params['tstop']
    for go_through in [1, 2]:
        if go_through==2:
            dpls=dpls_default
            net=net_default
        #--------------------
        # 1: simulation
        #--------------------
        if go_through==1:
            ax1=fig.add_subplot(321)
        else:
            ax1=fig.add_subplot(322)
        trial=0
        keep_for_mean=np.zeros((len(dpls),len(dpls[0].times)))
        for dpl in dpls:
            #ax1.plot(dpl.t,dpl.dpl['agg'],color=[.5, .5, .5],linewidth=1)
            keep_for_mean[trial,:]=dpl.data['agg']
            trial=trial+1      
        ax1.plot(dpl.times,np.mean(keep_for_mean,0),color='k',linewidth=3)
        if go_through==1:
            ax1.title.set_text(name)
        else:
            if supra:
                ax1.title.set_text('supra')
            else:
                ax1.title.set_text('default')
        ax1.set_xlim((0,xlim))
        ax1.legend(('Model','Data'))
        #ax1.legend(('Model','Contra','Ipsi'))
        ax1.set_ylim(ylms)
        #mark inputs
        these_params=net.params
        for temps in these_params:
            if temps[0:2]=='t_':
                if temps[4]=='p':
                    Colour='r'
                else:
                    Colour='g'
                ax1.arrow(these_params[temps],ylms[0],0,(ylms[1]-ylms[0])/10,color=Colour)                   
        #--------------------
        # 2: spiking    
        #--------------------
        if go_through==1:
            ax2=fig.add_subplot(323)
        else:
            ax2=fig.add_subplot(324)
        #spikes = np.array(sum(net.spiketimes, []))
        #gids = np.array(sum(net.spikegids, []))
        spikes = np.array((net.spikes._times[0]))
        gids = np.array((net.spikes._gids[0]))
        cell_types = ['L5_pyramidal', 'L5_basket', 'L2_pyramidal', 'L2_basket']
        cell_colours=['r','b','g',[.5, .5, .5]]
        count=0
        for cell in cell_types:
            these_spikes=[]
            these_times=[]
            for i in range(0,len(gids)):
                if gids[i]>= net.gid_dict[cell][0] and gids[i]<= net.gid_dict[cell][-1] :
                    #if any(gids[i] in s for s in net.gid_dict[cell]):
                    these_spikes.append(gids[i])
                    these_times.append(spikes[i])
            plt.scatter(these_times,these_spikes,color=cell_colours[count])
            count+=1
        ax2.legend(cell_types)
        ax2.get_yaxis().set_visible(False)
        ax2.title.set_text('Spiking')
        ax2.set_xlim((0,xlim))
        #--------------------
        # 3: dipoles
        #--------------------
        if go_through==1:
            ax3=fig.add_subplot(325)
        else:
            ax3=fig.add_subplot(326)
        trial=0
        keep_for_mean2=np.zeros((len(dpls),len(dpls[0].times)))
        keep_for_mean5=np.zeros((len(dpls),len(dpls[0].times)))
        for dpl in dpls:
            plt.plot(dpl.times,dpl.data['L2'],color='g',linewidth=1)
            plt.plot(dpl.times,dpl.data['L5'],color='r',linewidth=1)
            keep_for_mean2[trial,:]=dpl.data['L2']
            keep_for_mean5[trial,:]=dpl.data['L5']
            trial+=1
        leg1=plt.plot(dpl.times,np.mean(keep_for_mean2,0),color='g',linewidth=3)
        leg1=plt.plot(dpl.times,np.mean(keep_for_mean5,0),color='r',linewidth=3)
        plt.legend(('L2/3','L5'))
        ax3.title.set_text('Dipoles')
        ax3.set_xlim((0,xlim))
        ax3.set_ylim(ylms)
    plt.pause(.1)
    plt.savefig(dump_dir+'\\temp.png')
    img_path = dump_dir+'\\temp.png'
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
    plt.close()
prs.save(dump_dir+'\\new_pred_supra.pptx')